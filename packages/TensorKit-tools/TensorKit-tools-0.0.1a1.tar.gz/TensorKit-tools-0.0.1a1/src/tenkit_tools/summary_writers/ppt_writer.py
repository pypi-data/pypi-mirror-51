import csv
import os
import argparse
from pathlib import Path

import pptx
from pptx.util import Pt, Cm
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from ..utils import load_summary


TITLE_ONLY_SLIDE = 5
BLANK_SLIDE = 6
SLIDE_WIDTH = 9144000
SLIDE_HEIGHT = 5143500
TEMPLATE_NAME = str(Path(__file__).parent/'template.pptx')

FONT_NAME = 'Calibri'
FONT_SIZE = Pt(14)


def generate_table(slide, data_rows, column_names):
    # Setup table
    num_rows = len(data_rows)+1
    row_height = Cm(0.78)

    num_cols = 5
    col_widths = list(map(Cm, [3.15, 2.95, 4, 4.15, 3.5]))

    top, left = 476130, 2373459
    tablewidth = sum(col_widths)
    tableheight = (num_rows+1)*row_height

    tableshape = slide.shapes.add_table(
        num_rows, num_cols, left, top, tablewidth, tableheight
    )
    table = tableshape.table

    # Insert table content
    for cell_title, cell in zip(column_names, table.rows[0].cells):
        frame = cell.text_frame
        frame.clear()
        run = frame.paragraphs[0].add_run()
        run.text = cell_title

        font = run.font
        font.name = FONT_NAME
        font.size = FONT_SIZE
        font.bold = True

    for i, row_data in enumerate(data_rows):
        cellrow = table.rows[i+1]
        for cell_data, cell in zip(row_data.values(), cellrow.cells):
            frame = cell.text_frame
            frame.clear()
            run = frame.paragraphs[0].add_run()
            run.text = cell_data

            font = run.font
            font.name = FONT_NAME
            font.size = FONT_SIZE

    # Format table
    for col_width, column in zip(col_widths, table.columns):
        column.width = col_width


def generate_presentation(pres, data_rows, column_names, experiment_folder):
    # Setup slide
    pres.slide_width = SLIDE_WIDTH
    pres.slide_height = SLIDE_HEIGHT

    if len(pres.slides) > 0:
        slide = pres.slides[0]
    else:
        slide = pres.slides.add_slide(pres.slide_layouts[BLANK_SLIDE])
    
    generate_table(slide, data_rows, column_names)

    for experiment in sorted(experiment_folder.iterdir()):
        if not experiment.is_dir() or not (experiment/'summaries/summary.json').is_file():
            continue

        summary = load_summary(experiment)
        model = summary['model_type'].replace('_', ' ')
        for i, image in enumerate((experiment/'summaries'/'visualizations').iterdir()):
            slide = pres.slides.add_slide(pres.slide_layouts[TITLE_ONLY_SLIDE])
            slide.shapes.title.text = f'{model} model with {summary["model_rank"]} components'
            for paragraph in slide.shapes.title.text_frame.paragraphs:
            	paragraph.font.name = FONT_NAME
            	paragraph.font.bold = True
            	paragraph.font.size=Pt(18)
            	paragraph.alignment=PP_ALIGN.LEFT

            slide.shapes.title.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        
            image = slide.shapes.add_picture(str(image), Cm(i*2), Cm(i*2), height=Cm(5))
            image.left = int((SLIDE_WIDTH - image.width)/2)
            image.top = int((SLIDE_HEIGHT - image.height)/2)

    return pres


def create_ppt(parent_folder, csvpath='slide.csv', pptpath='summary.pptx'):
    if parent_folder is not None:
        csvpath = os.path.join(parent_folder, csvpath)
        pptpath = os.path.join(parent_folder, pptpath)

    with open(csvpath) as f:
        reader = csv.DictReader(f)
        data_rows = [row for row in reader]
        column_names = list(data_rows[0].keys())
    
    pres = pptx.Presentation(TEMPLATE_NAME)
    pres = generate_presentation(pres, data_rows, column_names, parent_folder)
    pres.save(pptpath)


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument("experiment_folder")
    parser.add_argument("csv_file")
    parser.add_argument("output_file")
    parser.add_argument("-p", "--parent_folder", default=None)
    parser.add_argument("-t", "--template", default="template.pptx")

    args = parser.parse_args()
    experiment_folder = Path(args.experiment_folder)
    csv_file = experiment_folder/args.csv_file
    output_file = args.output_file
    if args.parent_folder is not None:
        csv_file = os.path.join(args.parent_folder, csv_file)
        output_file = os.path.join(args.parent_folder, output_file)
    # Get slide data
    with open(csv_file) as f:
        reader = csv.DictReader(f)
        data_rows = [row for row in reader]
        column_names = list(data_rows[0].keys())

    pres = pptx.Presentation(args.template)
    pres = generate_presentation(pres, data_rows, column_names, experiment_folder)
    pres.save(output_file)
